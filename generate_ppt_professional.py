from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()

# Professional color palette
title_color = RGBColor(0x00, 0x8c, 0xd4)
subtitle_color = RGBColor(0x33, 0x4f, 0x6b)
body_color = RGBColor(0x10, 0x14, 0x29)
box_fill = RGBColor(0xf4, 0xf9, 0xfd)
box_border = RGBColor(0x00, 0x95, 0xd4)
text_color = RGBColor(0x1f, 0x2d, 0x3f)
accent_color = RGBColor(0x00, 0xac, 0xff)


def add_textbox(slide, left, top, width, height, text, font_size=Pt(18), bold=False, color=text_color, align=PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return textbox


def add_card(slide, left, top, width, height, title, lines):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = box_fill
    shape.line.color.rgb = box_border
    shape.line.width = Pt(1.5)

    text_frame = shape.text_frame
    text_frame.clear()
    title_p = text_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(18)
    title_p.font.bold = True
    title_p.font.color.rgb = accent_color
    title_p.space_after = Pt(8)

    for item in lines:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.level = 0
        p.space_after = Pt(4)

    return shape


def add_header(slide, title, subtitle):
    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.9), title, font_size=Pt(34), bold=True, color=title_color)
    add_textbox(slide, Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.5), subtitle, font_size=Pt(20), bold=False, color=subtitle_color)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.65), Inches(9.2), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()


def create_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, Inches(0.4), Inches(0.6), Inches(9.2), Inches(0.8), "MAX Enterprise AI Platform", font_size=Pt(44), bold=True, color=title_color)
    add_textbox(slide, Inches(0.4), Inches(1.6), Inches(9.2), Inches(0.6), "A clean, professional phase and architecture presentation.", font_size=Pt(20), bold=False, color=subtitle_color)
    add_textbox(slide, Inches(0.4), Inches(2.5), Inches(9.2), Inches(1.2), "Built for stakeholders, developers, and leadership reviews.", font_size=Pt(18), bold=False, color=text_color)


def create_vision_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Vision & Why This Matters", "AI observes, explains, and acts only with safe approval.")
    add_card(slide, Inches(0.4), Inches(2.4), Inches(4.5), Inches(3.2), "Why it matters", [
        "Reduce downtime with proactive monitoring",
        "Keep humans in control with approval and rollback",
        "Provide clear confidence before action",
    ])
    add_card(slide, Inches(5.1), Inches(2.4), Inches(4.5), Inches(3.2), "Business benefits", [
        "Faster issue resolution for teams",
        "Safer automation of critical workflows",
        "One platform for tech and business operations",
    ])


def create_architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Architecture Overview", "Modular system with UI, API, AI, and data layers.")
    add_card(slide, Inches(0.4), Inches(2.4), Inches(4.4), Inches(1.4), "Web & Mobile UI", [
        "Angular + PrimeNG responsive design",
        "Role-based screens and voice/chat interaction",
    ])
    add_card(slide, Inches(5.1), Inches(2.4), Inches(4.4), Inches(1.4), "Backend API", [
        "ASP.NET Core with JWT/RBAC",
        "Secure incident and action workflows",
    ])
    add_card(slide, Inches(0.4), Inches(4.1), Inches(4.4), Inches(1.4), "AI Orchestration", [
        "Python FastAPI + LangGraph",
        "LLM, STT, TTS, and root-cause analysis",
    ])
    add_card(slide, Inches(5.1), Inches(4.1), Inches(4.4), Inches(1.4), "Data Layer", [
        "Kafka streaming, PostgreSQL + pgVector",
        "Redis cache and real-time updates",
    ])


def create_phases_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Implementation Phases", "Three focused phases for build, intelligence, and scale.")
    add_card(slide, Inches(0.4), Inches(2.5), Inches(2.9), Inches(2.9), "Phase 1 — Foundation", [
        "Secure auth and OTP",
        "Role-based dashboard and JWT",
        "First app collector and Kafka pipeline",
        "Basic incident workflow",
    ])
    add_card(slide, Inches(3.5), Inches(2.5), Inches(2.9), Inches(2.9), "Phase 2 — AI Monitoring", [
        "Incident analysis and confidence score",
        "Approval, escalation, rollback",
        "Real-time alerts and voice support",
        "Business AI capabilities",
    ])
    add_card(slide, Inches(6.6), Inches(2.5), Inches(2.9), Inches(2.9), "Phase 3 — Autonomous AI", [
        "Agent orchestration and knowledge base",
        "Mobile-ready UX and broader integrations",
        "Safe automation with human gating",
    ])


def create_mobile_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Mobile Implementation", "Mobile-ready delivery is included in every phase.")
    add_card(slide, Inches(0.4), Inches(2.4), Inches(4.4), Inches(3.0), "Phase 1 — Mobile Foundation", [
        "Responsive UI and mobile-first layout design",
        "API contracts for compact mobile screens",
        "Voice/chat widget ready for phones and tablets",
    ])
    add_card(slide, Inches(5.1), Inches(2.4), Inches(4.4), Inches(3.0), "Phase 2 — Mobile UX", [
        "Touch-friendly approvals and alerts",
        "Compact incident summaries and charts",
        "Mobile-friendly notification and audio flows",
    ])
    add_card(slide, Inches(0.4), Inches(5.6), Inches(9.1), Inches(1.4), "Phase 3 — Mobile Scale", [
        "Optimize performance for phones/tablets",
        "Prepare backend APIs for native apps",
        "Mobile-ready escalation, rollback, and business actions",
    ])


def create_benefits_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Why Phase-Based Delivery", "A staged rollout reduces risk and clarifies progress.")
    add_card(slide, Inches(0.4), Inches(2.4), Inches(4.4), Inches(3.2), "Key Benefits", [
        "Clear milestones for teams and stakeholders",
        "Secure foundation before advanced automation",
        "Incremental validation and reduced risk",
        "Faster adoption with simpler reviews",
    ])
    add_card(slide, Inches(5.1), Inches(2.4), Inches(4.4), Inches(3.2), "Next Steps", [
        "Finalize Phase 1 MVP scope",
        "Build core auth, monitoring, and incident flow",
        "Deliver AI monitoring and business capabilities",
        "Prepare for mobile-ready Phase 3",
    ])

create_title_slide()
create_vision_slide()
create_architecture_slide()
create_phases_slide()
create_mobile_slide()
create_benefits_slide()

prs.save("MAX-AI-Platform-Phases-Professional-Updated.pptx")
